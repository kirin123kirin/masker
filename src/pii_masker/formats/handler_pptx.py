"""
PowerPointファイル (.pptx) ハンドラー
全スライドのテキストフレーム・表・ノートをマスクする。
書式（フォント・色・サイズ等）は保持する。
"""

import io
from pathlib import Path
from pptx import Presentation
from pptx.util import Pt
from pii_masker.engine.masker import Masker


def _mask_text_frame(tf, masker: Masker) -> tuple[list[str], list[str]]:
    """
    テキストフレームの各段落をマスクする。
    書式保持のため run 分配方式を使用。
    """
    originals, maskeds = [], []
    for para in tf.paragraphs:
        runs = para.runs
        if not runs:
            continue
        original = "".join(r.text for r in runs)
        if not original.strip():
            continue
        masked = masker.mask(original)
        originals.append(original)
        maskeds.append(masked)
        if original != masked:
            runs[0].text = masked
            for r in runs[1:]:
                r.text = ""
    return originals, maskeds


def _mask_table(table, masker: Masker) -> tuple[list[str], list[str]]:
    originals, maskeds = [], []
    for row in table.rows:
        for cell in row.cells:
            o_list, m_list = _mask_text_frame(cell.text_frame, masker)
            originals.extend(o_list)
            maskeds.extend(m_list)
    return originals, maskeds


def process_pptx(src: Path, masker: Masker) -> tuple[bytes, str, str]:
    """
    Returns:
        (output_bytes, original_text, masked_text)
    """
    prs = Presentation(str(src))

    all_original: list[str] = []
    all_masked: list[str] = []

    for slide in prs.slides:
        for shape in slide.shapes:
            # テキストフレーム
            if shape.has_text_frame:
                o_list, m_list = _mask_text_frame(shape.text_frame, masker)
                all_original.extend(o_list)
                all_masked.extend(m_list)

            # 表
            if shape.has_table:
                o_list, m_list = _mask_table(shape.table, masker)
                all_original.extend(o_list)
                all_masked.extend(m_list)

            # グループシェイプ（再帰）
            if shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
                for s in shape.shapes:
                    if s.has_text_frame:
                        o_list, m_list = _mask_text_frame(s.text_frame, masker)
                        all_original.extend(o_list)
                        all_masked.extend(m_list)

        # スライドノート
        if slide.has_notes_slide:
            notes_tf = slide.notes_slide.notes_text_frame
            o_list, m_list = _mask_text_frame(notes_tf, masker)
            all_original.extend(o_list)
            all_masked.extend(m_list)

    buf = io.BytesIO()
    prs.save(buf)
    output_bytes = buf.getvalue()

    return output_bytes, "\n".join(all_original), "\n".join(all_masked)
